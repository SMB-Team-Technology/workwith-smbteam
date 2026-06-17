"""
Audit PowerPoint — Farias Firm (David Farias)
SMB Team  |  Michael Kopp  |  June 17, 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ═══════════════════════════════════════════════════════════════════
# FILL — all values extracted from completed section files
# ═══════════════════════════════════════════════════════════════════

FIRM_NAME   = "Farias Firm (David Farias)"
SALES_REP   = "Michael Kopp"
OUTPUT_PATH = "farias-firm-david-farias/FariasFirm_June17_2026_Proposal.pptx"

WEBSITE_SCREENSHOT_PATH = None

# ── Slide 1 ──────────────────────────────────────────────────────
URGENCY_SCORE = "8"

PILLARS = [
    ("RED", "CRITICAL", "No website, GBP, or ads"),
    ("RED", "CRITICAL", "Owner handles all intake"),
    ("RED", "CRITICAL", "Solo — zero support staff"),
    ("RED", "CRITICAL", "No revenue target set"),
]
PILLAR_NAMES = ["Lead Generation", "Intake", "Team", "Profit Plan"]

FINDINGS = [
    ("neg", "Zero digital presence — invisible to every non-referral PI prospect in San Antonio"),
    ("neg", "David handles all intake and admin at 60–70 hrs/week — already at capacity"),
    ("neg", "Cash flow tied to settlement timing with no reserve, no revenue tracking"),
    ("pos", "Established TTLA network generated immediate cases at launch — real demand proven"),
]

COMPETITORS = [
    ("Carabin & Shaw",              "~1,207 reviews", "4.9 stars · owns the 3-pack"),
    ("George Salinas Injury",       "~414 reviews",   "4.9 stars · bilingual mktg"),
    ("Gamez Law Firm",              "~100 reviews",   "4.8 stars · 40+ yrs SA mkt"),
]
CLIENT_REVIEWS      = "0 reviews"
CLIENT_REVIEWS_NOTE = "← You are here"

STAGE_TEXT = "Stage 3: Solo Practitioner  →  Goal: Stage 6, Law Firm Owner"

# ── Slide 2 ──────────────────────────────────────────────────────
SLIDE_2_TITLE = "Your Growth Plan: 3 Priorities to Build a Firm on Systems"

SMB_MODEL_DESC = (
    "Four pillars must work together — Lead Generation, Intake, "
    "Team, and Profit. Missing any one means growth adds hours, not freedom. "
    "When all four are in place, David can focus exclusively on the "
    "courtroom work he built this firm to do."
)

GOAL_HEADLINE = "Referrals Only → Scalable Pipeline"
GOAL_DBM      = "Focus on trials, not on operations"

PRIORITIES = [
    (
        "Build the", "Marketing Engine", "1D4ED8",
        [
            "Build full PI website with attorney bio live",
            "Claim GBP and start closing the 1,200-review gap",
            "Launch Google Ads for PI keywords in San Antonio",
            "Activate Spanish-language market — 50-70% lower CPCs",
            "Enroll in LSA — appear above every competitor online",
        ],
    ),
    (
        "Fix Intake &", "Stop Losing Cases", "0F766E",
        [
            "Set up after-hours answering for PI calls 24/7",
            "Build call script to qualify and convert leads",
            "Install CRM so every lead source is tracked",
            "Measure conversion rate from contact to signed case",
            "Build intake that converts while David is in trial",
        ],
    ),
    (
        "Install Team &", "Profit Systems", "6D28D9",
        [
            "Set 12-month revenue target in first coach session",
            "Hire part-time assistant to free 15-20 hrs/week",
            "Track case costs vs settlements for profit clarity",
            "Build litigation cost reserve for settlement gaps",
            "Create role clarity so David focuses on litigation",
        ],
    ),
]

# ── Slide 3 ──────────────────────────────────────────────────────
PACKAGES = [
    (
        "FULL SERVICE MARKETING — STARTER",
        "$4,847", "$5,697/mo",
        "Website · Google Ads (EN+ES) · LSA · Meta · GBP · SEO",
        "1D4ED8",
    ),
    (
        "ELITE COACH",
        "$2,600", "$3,497/mo",
        "Weekly coaching · PI masterminds · Quarterly workshops",
        "6D28D9",
    ),
]

BUNDLE_TOTAL   = "$7,447 / mo"
BUNDLE_SAVINGS = "Save $1,747/mo by bundling"

AD_SPEND_NOTE = (
    "+ Recommended ad spend: $7,500–$25,000/mo paid directly to Google/Meta"
)

AVG_CASE_VALUE      = "$8,000"
CONSERVATIVE_LABEL  = "Conservative  (5–7 cases/mo):"
CONSERVATIVE_RESULT = "~$48,000 revenue · 6.4× ROAS"
AGGRESSIVE_LABEL    = "Aggressive  (22–24 cases/mo):"
AGGRESSIVE_RESULT   = "~$184,000 revenue · 7.4× ROAS"

TIMELINE = [
    ("Day 1",   "Ad accounts set up; GBP claim submitted"),
    ("Day 14",  "Website build begins; intake script deployed"),
    ("Week 2",  "GBP verified; review collection outreach starts"),
    ("Week 3",  "Ads go live on first PI landing page"),
    ("Month 3", "Coach kickoff: revenue target and 90-day plan set"),
]

CLOSING_QUOTE = (
    '"You built this firm to try cases — let us build the systems '
    'that make everything else run without you."'
)

# ═══════════════════════════════════════════════════════════════════
# LAYOUT ENGINE — DO NOT MODIFY BELOW THIS LINE
# ═══════════════════════════════════════════════════════════════════

def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

NAVY        = rgb("0D1F3C")
DARK_NAVY   = rgb("162947")
GOLD        = rgb("F5C400")
WHITE       = rgb("FFFFFF")
RED         = rgb("C0392B")
AMBER       = rgb("D97706")
GREEN       = rgb("16A34A")
SLATE       = rgb("64748B")
LIGHT_BLUE  = rgb("8BADD0")
FOOTER_TEXT = rgb("5A7A9A")
NEG_BG      = rgb("2C1010")
POS_BG      = rgb("0C2818")
NEG_TEXT    = rgb("F0BABA")
POS_TEXT    = rgb("86EFAC")
COMP_ALT    = rgb("F0F4FA")
ROI_BG      = rgb("ECFDF5")
ROI_GREEN   = rgb("065F46")
ROI_HILIGHT = rgb("D1FAE5")
NEAR_WHITE    = rgb("F0F4FA")
BUNDLE_SUB    = rgb("7CA0C0")
STRIKETHROUGH = rgb("334155")

STATUS_COLOR = {"RED": RED, "AMBER": AMBER, "GREEN": GREEN}
PRIORITY_LIGHT = {
    "1D4ED8": rgb("EFF6FF"),
    "0F766E": rgb("F0FDFA"),
    "6D28D9": rgb("F5F3FF"),
}

FONT = "Calibri"
LOGO_PATH = os.path.join(os.path.dirname(__file__), "smb_team_logo.png")

def emu(*inches):
    return tuple(int(x * 914400) for x in inches)


def add_rect(slide, left, top, w, h, fill=None, line=False):
    from pptx.util import Emu as E
    shape = slide.shapes.add_shape(1, E(int(left*914400)), E(int(top*914400)),
                                   E(int(w*914400)), E(int(h*914400)))
    shape.line.fill.background() if not line else None
    if not line:
        shape.line.color.rgb = WHITE
        shape.line.width = 0
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, left, top, w, h, size, color, bold=False,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    from pptx.util import Emu as E, Pt as P
    txb = slide.shapes.add_textbox(E(int(left*914400)), E(int(top*914400)),
                                   E(int(w*914400)), E(int(h*914400)))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = P(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    txb.text_frame.auto_size = None
    return txb


def add_image_or_placeholder(slide, path, left, top, w, h, label=""):
    from pptx.util import Emu as E
    if path and os.path.isfile(path):
        slide.shapes.add_picture(path, E(int(left*914400)), E(int(top*914400)),
                                 E(int(w*914400)), E(int(h*914400)))
    else:
        ph = add_rect(slide, left, top, w, h, fill=rgb("D1D5DB"))
        if label:
            add_text(slide, label, left+0.05, top+(h/2)-0.12, w-0.1, 0.24,
                     7, SLATE, align=PP_ALIGN.CENTER)


def add_footer(slide, page, total, logo_path):
    add_rect(slide, 0, 5.28, 10, 0.35, fill=NAVY)
    if logo_path and os.path.isfile(logo_path):
        from pptx.util import Emu as E
        slide.shapes.add_picture(logo_path,
                                 E(int(0.22*914400)), E(int(5.29*914400)),
                                 E(int(1.28*914400)), E(int(0.26*914400)))
    add_text(slide, f"CONFIDENTIAL  ·  Prepared by {SALES_REP} | SMB Team",
             1.76, 5.29, 5.80, 0.30, 8, FOOTER_TEXT)
    add_text(slide, f"{page} / {total}", 8.90, 5.29, 0.86, 0.30, 8, FOOTER_TEXT,
             align=PP_ALIGN.RIGHT)


def build_slide1(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_rect(slide, 0, 0, 10, 1.05, fill=NAVY)
    add_text(slide, f"LAW FIRM GROWTH AUDIT  ·  {FIRM_NAME.upper()}",
             0.32, 0.10, 9.40, 0.22, 8, GOLD, bold=True)
    add_text(slide, f"Where {FIRM_NAME} Stands Today",
             0.32, 0.34, 9.40, 0.60, 24, WHITE, bold=True)

    add_rect(slide, 0, 1.05, 0.16, 4.23, fill=GOLD)

    add_rect(slide, 0.28, 1.10, 1.42, 0.92, fill=RED)
    add_text(slide, URGENCY_SCORE, 0.28, 1.10, 1.42, 0.62, 28, WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "COMPETITIVE URGENCY", 0.28, 1.72, 1.42, 0.30, 6, WHITE, bold=True,
             align=PP_ALIGN.CENTER)

    pillar_xs = [1.84, 2.78, 3.72, 4.66]
    for i, (status, label, detail) in enumerate(PILLARS):
        x = pillar_xs[i]
        sc = STATUS_COLOR[status]
        add_rect(slide, x, 1.10, 0.88, 0.18, fill=sc)
        add_rect(slide, x, 1.28, 0.88, 0.74, fill=DARK_NAVY)
        add_text(slide, PILLAR_NAMES[i], x+0.06, 1.30, 0.76, 0.26, 8, WHITE, bold=True)
        add_text(slide, label, x+0.06, 1.55, 0.76, 0.20, 7, sc, bold=True)
        add_text(slide, detail, x+0.06, 1.74, 0.76, 0.26, 7, LIGHT_BLUE)

    add_text(slide, "KEY FINDINGS", 0.28, 2.14, 5.30, 0.24, 8, GOLD, bold=True)

    finding_ys = [2.42, 3.10, 3.78, 4.46]
    for i, (ftype, text) in enumerate(FINDINGS[:4]):
        y = finding_ys[i]
        bg   = NEG_BG if ftype == "neg" else POS_BG
        dot  = RED if ftype == "neg" else GREEN
        sym  = "✕" if ftype == "neg" else "✓"
        txt_color = NEG_TEXT if ftype == "neg" else POS_TEXT
        add_rect(slide, 0.28, y, 5.30, 0.62, fill=bg)
        add_rect(slide, 0.40, y+0.19, 0.24, 0.24, fill=dot)
        add_text(slide, sym, 0.40, y+0.17, 0.24, 0.26, 9, WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, text, 0.74, y+0.08, 4.76, 0.46, 10, txt_color)

    add_rect(slide, 5.75, 1.05, 4.25, 4.23, fill=WHITE)

    if WEBSITE_SCREENSHOT_PATH and os.path.isfile(WEBSITE_SCREENSHOT_PATH):
        from pptx.util import Emu as E
        slide.shapes.add_picture(WEBSITE_SCREENSHOT_PATH,
                                 E(int(5.82*914400)), E(int(1.10*914400)),
                                 E(int(4.10*914400)), E(int(2.00*914400)))

    add_rect(slide, 5.75, 1.12, 4.25, 0.76, fill=NAVY)
    add_rect(slide, 5.75, 1.12, 0.14, 0.76, fill=GOLD)
    add_text(slide, "YOU ARE HERE", 6.00, 1.14, 3.80, 0.22, 8, GOLD, bold=True)
    add_text(slide, STAGE_TEXT, 6.00, 1.36, 3.80, 0.44, 10, WHITE, bold=True)

    add_text(slide, "COMPETITOR LANDSCAPE", 5.90, 2.10, 3.90, 0.24, 8, SLATE, bold=True)

    comp_ys = [2.38, 2.76, 3.14]
    for i, (name, reviews, detail) in enumerate(COMPETITORS[:3]):
        y = comp_ys[i]
        bg = COMP_ALT if i % 2 == 0 else WHITE
        add_rect(slide, 5.75, y, 4.25, 0.34, fill=bg)
        add_text(slide, name, 5.92, y+0.05, 2.00, 0.24, 8, rgb("1E293B"))
        add_text(slide, reviews, 7.94, y+0.05, 1.00, 0.24, 8, GREEN)
        add_text(slide, detail, 8.96, y+0.07, 0.90, 0.22, 6, SLATE)

    add_rect(slide, 5.75, 3.52, 4.25, 0.34, fill=rgb("FFF0F0"))
    add_rect(slide, 5.75, 3.52, 0.10, 0.34, fill=RED)
    add_text(slide, FIRM_NAME, 5.92, 3.57, 2.00, 0.24, 9, RED, bold=True)
    add_text(slide, CLIENT_REVIEWS, 7.94, 3.57, 1.00, 0.24, 8, RED, bold=True)
    add_text(slide, CLIENT_REVIEWS_NOTE, 8.96, 3.59, 0.90, 0.22, 6, SLATE)

    add_footer(slide, 1, 3, LOGO_PATH)


def build_slide2(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_rect(slide, 0, 0, 10, 1.05, fill=NAVY)
    add_text(slide, f"LAW FIRM GROWTH AUDIT  ·  {FIRM_NAME.upper()}",
             0.32, 0.10, 9.40, 0.22, 8, GOLD, bold=True)
    add_text(slide, SLIDE_2_TITLE, 0.32, 0.34, 9.40, 0.60, 22, WHITE, bold=True)

    add_rect(slide, 0.20, 1.12, 2.72, 3.98, fill=NAVY)
    add_rect(slide, 0.20, 1.14, 2.72, 0.02, fill=GOLD)
    add_text(slide, "THE SMB TEAM MODEL", 0.32, 1.20, 2.48, 0.22, 8, GOLD, bold=True)
    add_text(slide, SMB_MODEL_DESC, 0.32, 1.46, 2.48, 1.72, 9, rgb("A8BFDA"))
    add_rect(slide, 0.20, 3.28, 2.72, 0.82, fill=GOLD)
    add_text(slide, GOAL_HEADLINE, 0.32, 3.30, 2.50, 0.30, 12, NAVY, bold=True)
    add_text(slide, GOAL_DBM, 0.32, 3.60, 2.50, 0.44, 9, NAVY)

    priority_xs = [3.08, 5.40, 7.72]
    bullet_ys   = [2.04, 2.70, 3.36, 4.02, 4.68]

    for col_i, (line1, line2, hex_color, bullets) in enumerate(PRIORITIES):
        x = priority_xs[col_i]
        ac = rgb(hex_color)
        light = PRIORITY_LIGHT.get(hex_color, rgb("F8F8FF"))

        add_rect(slide, x, 1.12, 2.24, 0.92, fill=ac)
        add_text(slide, f"0{col_i+1}", x+0.12, 1.14, 0.50, 0.30, 10, WHITE, bold=True)
        add_text(slide, line1, x+0.12, 1.44, 2.04, 0.28, 12, WHITE, bold=True)
        add_text(slide, line2, x+0.12, 1.70, 2.04, 0.28, 12, WHITE, bold=True)

        for row_i, bullet in enumerate(bullets[:5]):
            y = bullet_ys[row_i]
            bg = light if row_i % 2 == 0 else WHITE
            add_rect(slide, x, y, 2.24, 0.62, fill=bg)
            add_rect(slide, x+0.10, y+0.24, 0.10, 0.10, fill=ac)
            add_text(slide, bullet, x+0.26, y+0.06, 1.92, 0.52, 8, rgb("1E293B"))

    add_footer(slide, 2, 3, LOGO_PATH)


def build_slide3(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_rect(slide, 0, 0, 10, 1.05, fill=NAVY)
    add_text(slide, f"LAW FIRM GROWTH AUDIT  ·  {FIRM_NAME.upper()}",
             0.32, 0.10, 9.40, 0.22, 8, GOLD, bold=True)
    add_text(slide, "Your Investment & What Happens Next",
             0.32, 0.34, 9.40, 0.60, 22, WHITE, bold=True)

    pkg_ys = [1.12, 2.34]
    for i, (label, price, retail, services, hex_c) in enumerate(PACKAGES[:2]):
        y = pkg_ys[i]
        ac = rgb(hex_c)
        add_rect(slide, 0.22, y, 4.52, 1.12, fill=WHITE)
        add_rect(slide, 0.22, y, 0.20, 1.12, fill=ac)
        add_text(slide, label, 0.52, y+0.10, 4.16, 0.22, 8, ac, bold=True)
        add_text(slide, price, 0.52, y+0.30, 2.00, 0.48, 32, NAVY, bold=True)
        add_text(slide, "/mo", 2.04, y+0.42, 0.46, 0.28, 11, SLATE)
        add_text(slide, retail, 2.54, y+0.44, 1.00, 0.26, 11, STRIKETHROUGH)
        add_rect(slide, 2.54, y+0.55, 0.88, 0.01, fill=STRIKETHROUGH)
        add_text(slide, services, 0.52, y+0.84, 4.16, 0.22, 8, SLATE)

    add_rect(slide, 0.22, 3.56, 4.52, 0.72, fill=NAVY)
    add_text(slide, "BUNDLE TOTAL", 0.42, 3.60, 1.80, 0.26, 8, BUNDLE_SUB, bold=True)
    add_text(slide, BUNDLE_TOTAL, 0.42, 3.82, 2.40, 0.38, 22, GOLD, bold=True)
    add_text(slide, BUNDLE_SAVINGS, 2.92, 3.82, 1.90, 0.38, 8, BUNDLE_SUB)

    add_rect(slide, 0.22, 4.34, 4.52, 0.44, fill=rgb("EEF2F8"))
    add_text(slide, AD_SPEND_NOTE, 0.36, 4.37, 4.30, 0.38, 8, SLATE)

    add_rect(slide, 4.96, 1.12, 4.82, 1.44, fill=ROI_BG)
    add_text(slide, "PROJECTED RETURN ON AD SPEND",
             5.14, 1.18, 4.52, 0.24, 8, ROI_GREEN, bold=True)
    add_text(slide, "Average case value:", 5.14, 1.46, 2.10, 0.28, 9, rgb("1E293B"), bold=True)
    add_text(slide, AVG_CASE_VALUE, 7.36, 1.46, 2.30, 0.28, 9, ROI_GREEN)
    add_rect(slide, 5.08, 1.76, 4.58, 0.34, fill=ROI_HILIGHT)
    add_text(slide, CONSERVATIVE_LABEL, 5.14, 1.80, 2.10, 0.28, 9, rgb("1E293B"), bold=True)
    add_text(slide, CONSERVATIVE_RESULT, 7.36, 1.80, 2.30, 0.28, 9, ROI_GREEN, bold=True)
    add_text(slide, AGGRESSIVE_LABEL, 5.14, 2.14, 2.10, 0.28, 9, rgb("1E293B"), bold=True)
    add_text(slide, AGGRESSIVE_RESULT, 7.36, 2.14, 2.30, 0.28, 9, ROI_GREEN, bold=True)

    add_text(slide, "WHAT HAPPENS IN THE FIRST 90 DAYS",
             4.96, 2.70, 4.82, 0.26, 8, NAVY, bold=True)
    timeline_ys = [2.98, 3.37, 3.76, 4.15, 4.54]
    for i, (milestone, action) in enumerate(TIMELINE[:5]):
        y = timeline_ys[i]
        add_rect(slide, 5.02, y, 0.28, 0.28, fill=NAVY)
        add_text(slide, milestone, 5.40, y+0.01, 0.88, 0.28, 8, NAVY, bold=True)
        add_text(slide, action, 6.36, y+0.01, 3.34, 0.28, 8, rgb("1E293B"))

    add_rect(slide, 0, 4.84, 10, 0.44, fill=NAVY)
    add_text(slide, CLOSING_QUOTE, 0.30, 4.84, 9.40, 0.44, 9, LIGHT_BLUE, italic=True,
             align=PP_ALIGN.CENTER)

    add_footer(slide, 3, 3, LOGO_PATH)


# ── Assemble ──────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

build_slide1(prs)
build_slide2(prs)
build_slide3(prs)

os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True) if os.path.dirname(OUTPUT_PATH) else None
prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}  ({len(prs.slides)} slides)")
